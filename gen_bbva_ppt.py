"""
Genera una PPT editable que replica la Pantalla 3 (Panel BBVA)
de la presentación web presentacion-negocio-maduro.html
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colores (fieles al CSS) ──────────────────────────────────
BBVA_BLUE    = RGBColor(0x00, 0x44, 0x81)
BBVA_BLUE_LT = RGBColor(0xE8, 0xF0, 0xF8)
ORANGE       = RGBColor(0xEC, 0x68, 0x3E)
ORANGE_LT    = RGBColor(0xFD, 0xF0, 0xEB)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x10, 0x0D, 0x25)
GREY_100     = RGBColor(0xE8, 0xE8, 0xEE)
GREY_300     = RGBColor(0xB0, 0xB0, 0xBE)
GREY_400     = RGBColor(0x8E, 0x8E, 0xA0)
GREY_500     = RGBColor(0x6E, 0x6E, 0x80)
GREY_700     = RGBColor(0x3E, 0x3E, 0x4E)
BLUE_MED     = RGBColor(0x00, 0x44, 0x81)  # 45% opacity simulated
BLUE_MED_SIM = RGBColor(0x73, 0x95, 0xB8)
ORANGE_MED   = RGBColor(0xF4, 0xA8, 0x8E)

# ── Fuentes ──────────────────────────────────────────────────
FONT_SORA = 'Sora'
FONT_INST = 'Instrument Sans'
# Fallbacks si no están instaladas
FONT_TITLE = FONT_SORA
FONT_BODY  = FONT_INST

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

W = prs.slide_width
H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=Pt(1), corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    if corner_radius is not None:
        # Adjustments list for rounded corners
        if shape.adjustments:
            try:
                shape.adjustments[0] = corner_radius
            except:
                pass
    return shape

def add_text_box(slide, left, top, width, height, text, font_name=FONT_BODY, font_size=Pt(12),
                 font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_rich_text(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    """runs = list of dicts: {text, font, size, color, bold}"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    p = tf.paragraphs[0]
    if spacing:
        p.space_after = spacing
    p.alignment = alignment
    for i, r in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = r['text']
        else:
            run = p.add_run()
            run.text = r['text']
        run.font.name = r.get('font', FONT_BODY)
        run.font.size = r.get('size', Pt(12))
        run.font.color.rgb = r.get('color', BLACK)
        run.font.bold = r.get('bold', False)
    return txBox

def add_multiline_box(slide, left, top, width, height, paragraphs_data, alignment=PP_ALIGN.LEFT):
    """paragraphs_data = list of list of run-dicts"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for pi, runs in enumerate(paragraphs_data):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        for ri, r in enumerate(runs):
            run = p.add_run()
            run.text = r['text']
            run.font.name = r.get('font', FONT_BODY)
            run.font.size = r.get('size', Pt(12))
            run.font.color.rgb = r.get('color', BLACK)
            run.font.bold = r.get('bold', False)
        if 'space_after' in runs[0]:
            p.space_after = runs[0]['space_after']
        if 'space_before' in runs[0]:
            p.space_before = runs[0]['space_before']
    return txBox


# ══════════════════════════════════════════════════════════════
# BACKGROUND — fondo oscuro como la presentación
# ══════════════════════════════════════════════════════════════
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x10, 0x0D, 0x25)

# ══════════════════════════════════════════════════════════════
# PANEL BLANCO — el panel que sube (96% de altura)
# ══════════════════════════════════════════════════════════════
panel_margin_top = Inches(0.3)
panel_left = 0
panel_top = panel_margin_top
panel_w = W
panel_h = H - panel_margin_top

panel_bg = add_rect(slide, panel_left, panel_top, panel_w, panel_h,
                    fill=WHITE, corner_radius=0.02)

# ══════════════════════════════════════════════════════════════
# CABECERA AZUL BBVA
# ══════════════════════════════════════════════════════════════
header_h = Inches(0.85)
header = add_rect(slide, panel_left, panel_top, panel_w, header_h,
                  fill=BBVA_BLUE, corner_radius=0.015)

# Nombre cliente
add_text_box(slide, Inches(0.5), panel_top + Inches(0.12), Inches(4), Inches(0.45),
             "BBVA", FONT_TITLE, Pt(32), WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# Tagline
add_text_box(slide, Inches(0.5), panel_top + Inches(0.52), Inches(4), Inches(0.25),
             "CIB · CF&Banking · Geografías", FONT_BODY, Pt(11), RGBColor(0xCC, 0xDD, 0xEE),
             bold=False)

# Botón cerrar (✕)
add_text_box(slide, W - Inches(0.7), panel_top + Inches(0.2), Inches(0.4), Inches(0.4),
             "✕", FONT_BODY, Pt(16), RGBColor(0xAA, 0xCC, 0xEE), bold=False,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════
# EVOLUCIÓN FY — banda con borde superior azul
# ══════════════════════════════════════════════════════════════
evo_top = panel_top + header_h
evo_h = Inches(0.95)

# Línea azul superior
add_rect(slide, panel_left, evo_top, panel_w, Pt(3), fill=BBVA_BLUE)

evolution = [
    {"fy": "FY25", "rev": "29", "ftes": "400", "obj": False, "growth": None},
    {"fy": "FY26", "rev": "38", "ftes": "548", "obj": False, "growth": None},
    {"fy": "FY27", "rev": "47", "ftes": "685", "obj": True,  "growth": "↑ +60% acumulado"},
]

evo_col_w = W / 3
for i, e in enumerate(evolution):
    col_left = evo_col_w * i
    col_top = evo_top + Pt(6)

    # FY label
    fy_label = e["fy"]
    if e["obj"]:
        fy_label += " obj."
    add_text_box(slide, col_left + Inches(0.3), col_top + Inches(0.02), Inches(2), Inches(0.22),
                 fy_label, FONT_TITLE, Pt(10), GREY_500, bold=False)

    # Revenue
    add_rich_text(slide, col_left + Inches(0.3), col_top + Inches(0.2), Inches(2.5), Inches(0.35),
                  [
                      {'text': e["rev"], 'font': FONT_TITLE, 'size': Pt(26), 'color': BLACK, 'bold': True},
                      {'text': ' €MM', 'font': FONT_TITLE, 'size': Pt(10), 'color': GREY_400, 'bold': False}
                  ])

    # FTEs
    add_text_box(slide, col_left + Inches(0.3), col_top + Inches(0.52), Inches(2), Inches(0.2),
                 f"{e['ftes']} Ftes", FONT_TITLE, Pt(12), BLACK, bold=True)

    # Growth badge
    if e["growth"]:
        badge = add_rect(slide, col_left + Inches(0.3), col_top + Inches(0.72), Inches(1.8), Inches(0.2),
                         fill=BBVA_BLUE_LT, corner_radius=0.5)
        add_text_box(slide, col_left + Inches(0.3), col_top + Inches(0.72), Inches(1.8), Inches(0.2),
                     e["growth"], FONT_TITLE, Pt(9), BBVA_BLUE, bold=True,
                     alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Separador vertical
    if i < 2:
        sep_left = col_left + evo_col_w - Pt(1)
        add_rect(slide, sep_left, evo_top + Inches(0.08), Pt(1.5), evo_h - Inches(0.16), fill=GREY_100)

# ══════════════════════════════════════════════════════════════
# DOS COLUMNAS: ESTADO ACTUAL (izq) + RETOS (der)
# ══════════════════════════════════════════════════════════════
cols_top = evo_top + evo_h + Inches(0.05)
cols_h = Inches(3.6)
col_w = W / 2

# Separador vertical central
add_rect(slide, col_w - Pt(0.75), cols_top, Pt(1.5), cols_h, fill=GREY_100)

# ── ESTADO ACTUAL (columna izquierda) ────────────────────────
sa_left = Inches(0.5)
sa_top = cols_top + Inches(0.1)

# Icono + título
# Icono circular
icon_bg = add_rect(slide, sa_left, sa_top, Inches(0.35), Inches(0.35),
                   fill=BBVA_BLUE_LT, corner_radius=0.3)
add_text_box(slide, sa_left + Inches(0.03), sa_top + Inches(0.02), Inches(0.3), Inches(0.3),
             "✓", FONT_TITLE, Pt(14), BBVA_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text_box(slide, sa_left + Inches(0.45), sa_top + Inches(0.02), Inches(3), Inches(0.3),
             "Estado Actual", FONT_TITLE, Pt(16), BBVA_BLUE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

# Cards de Estado Actual
estado_actual = [
    {"badge": "Especialista", "title": "PERCIBIDOS COMO ESPECIALISTA",
     "body": "Aunque concentrada la expertise entre Ingeniería CIB, Riesgos Global y Finanzas.",
     "impact": "high"},
    {"badge": "Posicionamiento", "title": "BUEN POSICIONAMIENTO PARA ESCALAR",
     "body": "Firmado contrato POST 2.0 para la práctica en Ingeniería Global, FARO para Riesgos y abiertos todos los necesarios para Geografías.",
     "impact": "high"},
    {"badge": "Personas", "title": "BUENA ESTRUCTURA INICIAL DE PERSONAS",
     "body": "Con Conocimiento y Expertise para ir asumiendo nuevas responsabilidades en el crecimiento y nueva posición a coger.",
     "impact": "medium"},
]

card_top = sa_top + Inches(0.55)
card_w = col_w - Inches(1.1)

for j, item in enumerate(estado_actual):
    ct = card_top + j * Inches(1.05)
    bar_color = BBVA_BLUE if item["impact"] == "high" else BLUE_MED_SIM

    # Card border
    card = add_rect(slide, sa_left, ct, card_w, Inches(0.92),
                    fill=WHITE, border_color=GREY_100, border_width=Pt(1.5), corner_radius=0.02)

    # Left accent bar
    add_rect(slide, sa_left, ct + Inches(0.08), Pt(4), Inches(0.76), fill=bar_color)

    # Badge
    add_text_box(slide, sa_left + Inches(0.18), ct + Inches(0.06), card_w - Inches(0.3), Inches(0.18),
                 item["badge"], FONT_BODY, Pt(8), BBVA_BLUE, bold=True)

    # Title
    add_text_box(slide, sa_left + Inches(0.18), ct + Inches(0.24), card_w - Inches(0.3), Inches(0.25),
                 item["title"], FONT_TITLE, Pt(11), BLACK, bold=True)

    # Body
    add_text_box(slide, sa_left + Inches(0.18), ct + Inches(0.48), card_w - Inches(0.3), Inches(0.42),
                 item["body"], FONT_BODY, Pt(10), GREY_700, bold=False)


# ── RETOS (columna derecha) ──────────────────────────────────
rt_left = col_w + Inches(0.4)
rt_top = cols_top + Inches(0.1)

# Icono + título
icon_bg2 = add_rect(slide, rt_left, rt_top, Inches(0.35), Inches(0.35),
                    fill=ORANGE_LT, corner_radius=0.3)
add_text_box(slide, rt_left + Inches(0.03), rt_top + Inches(0.02), Inches(0.3), Inches(0.3),
             "⚡", FONT_TITLE, Pt(14), ORANGE, bold=True,
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text_box(slide, rt_left + Inches(0.45), rt_top + Inches(0.02), Inches(3), Inches(0.3),
             "Retos", FONT_TITLE, Pt(16), ORANGE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)

retos = [
    {"badge": "Ingeniería Global", "title": "LIDERAR INGENIERÍA DE MANERA GLOBAL",
     "body": "Llegar a ser proveedor en Ingeniería con relevancia de manera Global, sin perder la Esencia de Especialistas en CIB y manteniendo su liderazgo.",
     "impact": "high"},
    {"badge": "Geografías", "title": "SER RELEVANTES EN LAS GEOGRAFÍAS",
     "body": "Desde el Escalado en México a estar en Perú y Colombia, acompañando la apertura de Brasil.",
     "impact": "high"},
    {"badge": "Diversificación", "title": "DIVERSIFICAR RESPECTO DE CIB",
     "body": "Pasar del 80% a 60% sin dejar de crecer en CIB.",
     "impact": "medium"},
]

card_top_r = rt_top + Inches(0.55)
card_w_r = col_w - Inches(1.0)

for j, item in enumerate(retos):
    ct = card_top_r + j * Inches(1.05)
    bar_color = ORANGE if item["impact"] == "high" else ORANGE_MED

    # Card border
    card = add_rect(slide, rt_left, ct, card_w_r, Inches(0.92),
                    fill=WHITE, border_color=GREY_100, border_width=Pt(1.5), corner_radius=0.02)

    # Left accent bar
    add_rect(slide, rt_left, ct + Inches(0.08), Pt(4), Inches(0.76), fill=bar_color)

    # Badge
    add_text_box(slide, rt_left + Inches(0.18), ct + Inches(0.06), card_w_r - Inches(0.3), Inches(0.18),
                 item["badge"], FONT_BODY, Pt(8), ORANGE, bold=True)

    # Title
    add_text_box(slide, rt_left + Inches(0.18), ct + Inches(0.24), card_w_r - Inches(0.3), Inches(0.25),
                 item["title"], FONT_TITLE, Pt(11), BLACK, bold=True)

    # Body
    add_text_box(slide, rt_left + Inches(0.18), ct + Inches(0.48), card_w_r - Inches(0.3), Inches(0.42),
                 item["body"], FONT_BODY, Pt(10), GREY_700, bold=False)


# ══════════════════════════════════════════════════════════════
# ÁREAS — footer con botones
# ══════════════════════════════════════════════════════════════
areas_top = cols_top + cols_h + Inches(0.15)
areas_label_left = Inches(0.5)

add_text_box(slide, areas_label_left, areas_top, Inches(1), Inches(0.22),
             "ÁREAS", FONT_BODY, Pt(9), GREY_500, bold=True)

sections = [
    {"label": "CIB"},
    {"label": "CF & Banking"},
]

btn_left = areas_label_left
btn_top = areas_top + Inches(0.28)
for s in sections:
    btn_w = Inches(1.8)
    btn_h = Inches(0.42)
    btn = add_rect(slide, btn_left, btn_top, btn_w, btn_h,
                   fill=BBVA_BLUE, corner_radius=0.12)

    # Label + arrow
    add_rich_text(slide, btn_left, btn_top, btn_w, btn_h,
                  [
                      {'text': s["label"], 'font': FONT_BODY, 'size': Pt(11), 'color': WHITE, 'bold': True},
                      {'text': '  →', 'font': FONT_BODY, 'size': Pt(13), 'color': RGBColor(0xAA, 0xCC, 0xEE), 'bold': False}
                  ],
                  alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    btn_left += btn_w + Inches(0.2)


# ══════════════════════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════════════════════
output_path = r"c:\Users\angel.rodriguez\Documents\test\Proyectos_1\BBVA_Panel_Pantalla3.pptx"
prs.save(output_path)
print(f"PPT generada: {output_path}")
